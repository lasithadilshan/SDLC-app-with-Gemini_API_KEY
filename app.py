import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
import pptx
import pandas as pd
import os
import time
from langchain.chains.retrieval_qa.base import RetrievalQA
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI

st.set_page_config(
    page_title="SDLC Automate APP",
    page_icon="images/favicon.png"
)

# Get the API key from Streamlit secrets
GEMINI_API_KEY = st.secrets["GEMINI_API_KEY"]

# Streamlit sidebar setup
with st.sidebar:
    st.title("Your BRD Documents")
    uploaded_file = st.file_uploader("Upload a file to generate user stories", type=["pdf", "docx", "txt", "xlsx", "pptx"])

# Function to extract text from various file types
@st.cache_resource
def extract_text_from_file(file):
    """Extracts text based on file type, with caching for faster retrieval."""
    text = ""
    file_ext = os.path.splitext(file.name)[1].lower()

    # Handle PDF files
    if file_ext == ".pdf":
        pdf_reader = PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text()

    # Handle Word (.docx) files
    elif file_ext == ".docx":
        doc = Document(file)
        for para in doc.paragraphs:
            text += para.text + "\n"

    # Handle text (.txt) files
    elif file_ext == ".txt":
        text = file.read().decode("utf-8")

    # Handle Excel files (.xlsx, .xls)
    elif file_ext in [".xlsx", ".xls"]:
        df = pd.read_excel(file)
        text = df.to_string()

    # Handle PowerPoint files (.pptx, .ppt)
    elif file_ext in [".pptx", ".ppt"]:
        ppt = pptx.Presentation(file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    return text

# Process the uploaded file and extract text for the vector store
@st.cache_resource
def process_uploaded_file(uploaded_file):
    return extract_text_from_file(uploaded_file) if uploaded_file else ""

# Function to create vector store from extracted text
@st.cache_resource
def create_vector_store(text):
    text_splitter = RecursiveCharacterTextSplitter(
        separators="\n",
        chunk_size=800,
        chunk_overlap=50,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    embeddings = OpenAIEmbeddings(google_api_key=GEMINI_API_KEY)
    return FAISS.from_texts(chunks, embeddings)

# Streamlit app setup
st.header("BRD to User Story, Test Case, Cucumber Script, and Selenium Script")

# Set up tabs for different functionalities
tab1, tab2, tab3, tab4 = st.tabs(["User Story Generation", "User Story to Test Case", "Test Case to Cucumber Script", "Test Case to Selenium Script"])

# User Story Generation Tab
with tab1:
    start_time = time.time()
    if uploaded_file:
        text = process_uploaded_file(uploaded_file)
        if text:
            vector_store = create_vector_store(text)
            prompt_message = (
                "Think of yourself as a senior business analyst. Your responsibility is to read the Business Requirement Document "
                "and write the User Stories according to that BRD. Think step-by-step and write all possible user stories "
                "for the Business Requirement Document."
            )
            start_query_time = time.time()
            matches = vector_store.similarity_search(prompt_message, k=3)  # Retrieve top 3 similar texts
            llm = ChatGoogleGenerativeAI(
                google_api_key=GEMINI_API_KEY,
                temperature=0.9,
                max_output_tokens=500,
                model_name="gemini-pro"
            )
            qa_chain = RetrievalQA.from_chain_type(
                llm=llm,
                chain_type="stuff",
                retriever=vector_store.as_retriever()
            )
            response = qa_chain.invoke({"query": prompt_message})
            st.write(response['result'])

            # Display timing info for performance insights
            st.write(f"Document loading time: {time.time() - start_time:.2f} seconds")
            st.write(f"Query processing time: {time.time() - start_query_time:.2f} seconds")
    else:
        st.write("Please upload a BRD document in the sidebar to generate user stories.")

# User Story to Test Case Tab
with tab2:
    st.subheader("Convert User Story to Test Case")
    user_story_text = st.text_area("Enter the user story text here to generate test cases:")

    if st.button("Generate Test Cases"):
        if user_story_text:
            test_case_prompt = (
                    "Think of yourself as a senior QA engineer. Your responsibility is to read the user story provided and generate "
                    "all possible test cases. Think in a structured way, covering functional and edge cases where applicable. "
                    "Here is the user story: \n\n" + user_story_text
            )
            start_test_case_time = time.time()
            response = qa_chain.invoke({"query": test_case_prompt})
            st.write(response['result'])
            st.write(f"Test case generation time: {time.time() - start_test_case_time:.2f} seconds")
        else:
            st.write("Please enter a user story to generate test cases.")

# Test Case to Cucumber Script Tab
with tab3:
    st.subheader("Convert Test Case to Cucumber Script")
    test_case_text = st.text_area("Enter the test case text here to generate Cucumber script:")

    if st.button("Generate Cucumber Script"):
        if test_case_text:
            cucumber_prompt = (
                    "Think of yourself as a test automation engineer. Your task is to convert the following test case into a Cucumber "
                    "script using Gherkin syntax. Make sure to include all scenarios with Given, When, Then steps as applicable. "
                    "Here is the test case: \n\n" + test_case_text
            )
            start_cucumber_time = time.time()
            response = qa_chain.invoke({"query": cucumber_prompt})
            st.write(response['result'])
            st.write(f"Cucumber script generation time: {time.time() - start_cucumber_time:.2f} seconds")
        else:
            st.write("Please enter a test case to generate a Cucumber script.")

# Test Case to Selenium Script Tab
with tab4:
    st.subheader("Convert Test Case to Selenium Script")
    selenium_test_case_text = st.text_area("Enter the test case text here to generate Selenium script:")

    if st.button("Generate Selenium Script"):
        if selenium_test_case_text:
            selenium_prompt = (
                    "Assume you are a test automation engineer specializing in Selenium. Your task is to convert the following test case "
                    "into a Selenium WebDriver script using Python. Ensure to include all steps to perform the actions in the test case, "
                    "such as locating elements, interacting with the web page, and validating outcomes. Here is the test case: \n\n" + selenium_test_case_text
            )
            start_selenium_time = time.time()
            response = qa_chain.invoke({"query": selenium_prompt})
            st.write(response['result'])
            st.write(f"Selenium script generation time: {time.time() - start_selenium_time:.2f} seconds")
        else:
            st.write("Please enter a test case to generate a Selenium script.")
